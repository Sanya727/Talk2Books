import os
import fitz
from docx import Document
from pptx import Presentation

from langchain_core.documents import Document as LCDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter


def extract_pdf(path):

    doc = fitz.open(path)

    text = ""

    for page in doc:
        text += page.get_text()

    return text


def extract_docx(path):

    doc = Document(path)

    text = ""

    for p in doc.paragraphs:
        text += p.text + "\n"

    return text


def extract_ppt(path):

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def load_documents(folder):

    docs = []

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100
    )

    for file in os.listdir(folder):

        path = os.path.join(folder, file)

        if file.endswith(".pdf"):
            text = extract_pdf(path)

        elif file.endswith(".docx"):
            text = extract_docx(path)

        elif file.endswith(".pptx"):
            text = extract_ppt(path)

        elif file.endswith(".txt"):
           with open(path, encoding="utf-8", errors="ignore") as f:
            text = f.read()
        else:
            continue

        chunks = splitter.split_text(text)

        for c in chunks:

            docs.append(
                LCDocument(
                    page_content=c,
                    metadata={"source": file}
                )
            )

    return docs


def load_documents_from_text(text, source):

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100
    )

    chunks = splitter.split_text(text)

    docs = []

    for c in chunks:

        docs.append(
            LCDocument(
                page_content=c,
                metadata={"source": source}
            )
        )

    return docs